import io
import os
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor

# ✅ Correct folder name (make sure this matches your actual folder)
TEMPLATE_PATH = "templates/sample_ppt_template.pptx"


# ------------------------------------------------------------
# Placeholder Replacement Engine (Improved)
# ------------------------------------------------------------
def fill_placeholders(shape, data):
    """
    Recursively search and replace placeholders.
    Turns 'Information Missing' RED.
    Handles whitespace safely.
    """

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text

                for key, value in data.items():
                    placeholder = f"{{{key}}}"

                    # Safe check (strip prevents whitespace mismatch issues)
                    if placeholder in original_text.strip():

                        # Normalize value safely
                        safe_value = "" if value is None else str(value).strip()

                        # Missing data logic
                        if not safe_value or safe_value.lower() == "information missing":
                            run.text = original_text.replace(
                                placeholder, "Information Missing"
                            )
                            run.font.color.rgb = RGBColor(255, 0, 0)  # Red
                            run.font.bold = True
                        else:
                            run.text = original_text.replace(
                                placeholder, safe_value
                            )

    # Recursive table handling
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                fill_placeholders(cell, data)


# ------------------------------------------------------------
# Mapping Layer (Hardened)
# ------------------------------------------------------------
def map_resume_to_template(candidate_data):
    """
    Maps candidate data to flat keys used in PPT slides 1–5.
    Handles None values safely.
    """

    candidate_data = candidate_data or {}

    # Slide 1
    flat_data = {
        "FULL_NAME": candidate_data.get("name") or "",
        "CURRENT_ROLE": candidate_data.get("current_role") or "",
        "PROFESSIONAL_SUMMARY": candidate_data.get("summary") or "",
        "TECHNICAL_SKILLS": (
            ", ".join(candidate_data.get("tech_stack", []))
            if isinstance(candidate_data.get("tech_stack"), list)
            else candidate_data.get("tech_stack") or ""
        ),
        "EDUCATION_DETAILS": candidate_data.get("education") or "",
    }

    # Slides 2–5 (Projects)
    projects = candidate_data.get("work_experience", []) or []

    for i in range(1, 5):
        proj = projects[i - 1] if len(projects) >= i else {}

        name = proj.get("project_title") or proj.get("company") or ""
        duration = proj.get("dates") or ""
        role = proj.get("role") or candidate_data.get("current_role") or ""
        desc = proj.get("description") or ""

        bullets = proj.get("bullets", []) or []
        resp = "\n".join([f"• {b}" for b in bullets]) if bullets else ""

        flat_data.update({
            f"PROJECT{i}_NAME": name,
            f"DURATION_PROJECT{i}": duration,
            f"DURATION_project{i}": duration,  # case flexibility
            f"ROLE_PROJECT{i}": role,
            f"ROLE_project{i}": role,
            f"Project{i}_Description": desc,
            f"Responsibilities_Project{i}": resp,
            f"Responsibilities_project{i}": resp,
        })

    return flat_data


# ------------------------------------------------------------
# PPT Generator (Production Ready)
# ------------------------------------------------------------
def generate_candidate_ppt(candidate_data: dict) -> bytes | None:
    """
    Generates PPT and marks missing fields in red.
    """

    try:
        if not os.path.exists(TEMPLATE_PATH):
            st.error(f"Template not found at {TEMPLATE_PATH}")
            return None

        prs = Presentation(TEMPLATE_PATH)
        data_map = map_resume_to_template(candidate_data)

        for slide in prs.slides:
            for shape in slide.shapes:
                fill_placeholders(shape, data_map)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        return buf.read()

    except Exception as e:
        st.error(f"PPT Generation Error: {e}")
        return None